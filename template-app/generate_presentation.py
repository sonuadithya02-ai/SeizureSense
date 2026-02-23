#!/usr/bin/env python3
"""
Generate a beautiful minimalistic PowerPoint presentation for Seizure Detection Project
Clean design without overlapping elements
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette - Mellow but bright
COLORS = {
    'primary': RGBColor(78, 165, 175),      # Soft teal
    'secondary': RGBColor(245, 166, 128),   # Soft coral/peach
    'accent': RGBColor(147, 112, 219),      # Soft purple
    'dark': RGBColor(55, 65, 81),           # Dark gray
    'white': RGBColor(255, 255, 255),       # White
    'mint': RGBColor(110, 195, 162),        # Mint green
    'sky': RGBColor(135, 182, 217),         # Sky blue
    'warm': RGBColor(248, 180, 150),        # Warm peach
    'light_bg': RGBColor(248, 250, 252),    # Very light gray
}

def set_slide_background(slide, color):
    """Set solid background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, color):
    """Add a colored bar at top of slide"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=None, align='left'):
    """Add text box with specified formatting"""
    if color is None:
        color = COLORS['dark']
    
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    
    return shape

def add_rounded_box(slide, left, top, width, height, fill_color, text="", font_size=14, text_color=None):
    """Add rounded rectangle with text"""
    if text_color is None:
        text_color = COLORS['white']
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return shape

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, COLORS['light_bg'])

# Title bar
add_rounded_box(slide1, 0, 0, 13.333, 0.2, COLORS['primary'])

# Main title
add_text(slide1, "MULTI-SENSOR SEIZURE", 0.5, 2.3, 12.333, 0.8, font_size=48, bold=True, color=COLORS['dark'], align='center')
add_text(slide1, "DETECTION SYSTEM", 0.5, 3.1, 12.333, 0.8, font_size=48, bold=True, color=COLORS['primary'], align='center')

# Subtitle
add_text(slide1, "Using Machine Learning and Wearable Sensors", 0.5, 4.2, 12.333, 0.5, font_size=24, color=COLORS['dark'], align='center')

# Tech stack
add_text(slide1, "ESP32-S3  •  TensorFlow Lite  •  BLE Alert System", 0.5, 4.9, 12.333, 0.5, font_size=18, color=COLORS['accent'], align='center')

# Bottom info
add_text(slide1, "Final Year Project  •  2024-2025", 0.5, 6.5, 12.333, 0.4, font_size=14, color=COLORS['dark'], align='center')

# ============================================================================
# SLIDE 2: Problem Statement
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, COLORS['light_bg'])
add_title_bar(slide2, COLORS['primary'])

add_text(slide2, "Problem Statement", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['primary'], align='center')

# Problem description
add_text(slide2, "Epilepsy affects over 50 million people worldwide. Seizures can occur unexpectedly,\nleading to injuries, accidents, and Sudden Unexpected Death in Epilepsy (SUDEP).", 
         0.8, 1.2, 11.7, 0.9, font_size=16, color=COLORS['dark'], align='center')

# Statistics boxes
add_rounded_box(slide2, 1, 2.4, 3.5, 1.3, COLORS['primary'], "50+ Million\nAffected Globally", 18)
add_rounded_box(slide2, 5, 2.4, 3.5, 1.3, COLORS['secondary'], "1 in 26 People\nDevelop Epilepsy", 18)
add_rounded_box(slide2, 9, 2.4, 3.5, 1.3, COLORS['accent'], "SUDEP Risk:\n1 in 1000/Year", 18)

# Need for solution
add_text(slide2, "Need for Solution:", 0.8, 4.1, 5, 0.4, font_size=20, bold=True, color=COLORS['primary'])

needs = [
    "Real-time seizure detection for immediate caregiver alert",
    "Non-invasive, wearable solution for continuous monitoring", 
    "Multi-modal approach to reduce false alarms"
]
for i, need in enumerate(needs):
    add_text(slide2, f"✓  {need}", 1, 4.6 + i*0.5, 11, 0.4, font_size=15, color=COLORS['dark'])

# ============================================================================
# SLIDE 3: Literature Survey (Part 1)
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, COLORS['light_bg'])
add_title_bar(slide3, COLORS['secondary'])

add_text(slide3, "Literature Survey", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['secondary'], align='center')

# Table header
add_rounded_box(slide3, 0.3, 1.1, 4.2, 0.5, COLORS['dark'], "PAPER INFORMATION", 11)
add_rounded_box(slide3, 4.5, 1.1, 8.5, 0.5, COLORS['dark'], "FOCUS AREA, KEY CONTRIBUTIONS, LIMITATIONS", 11)

# Paper 1: Ulate-Campos et al.
add_rounded_box(slide3, 0.3, 1.65, 4.2, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper1_info = slide3.shapes.add_textbox(Inches(0.4), Inches(1.7), Inches(4), Inches(1.3))
tf = paper1_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "1. Ulate-Campos, Coughlin, Gaínza-Lein, et al."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"Automated seizure detection systems and their effectiveness"'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Year: 2016"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3, 4.5, 1.65, 8.5, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper1_details = slide3.shapes.add_textbox(Inches(4.6), Inches(1.7), Inches(8.3), Inches(1.3))
tf = paper1_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: Automated seizure detection across various seizure types"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: Comprehensive analysis categorizing systems by methodologies"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Limitation: Many devices prone to false alarms or missed detections"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['secondary']

# Paper 2: Selvi et al. (2023)
add_rounded_box(slide3, 0.3, 3.1, 4.2, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper2_info = slide3.shapes.add_textbox(Inches(0.4), Inches(3.15), Inches(4), Inches(1.3))
tf = paper2_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "2. Selvi, Thamil, et al."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"An Intelligent and Wearable Epileptic Seizure Detection Device"'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "IEEE ICCEBS, Year: 2023"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3, 4.5, 3.1, 8.5, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper2_details = slide3.shapes.add_textbox(Inches(4.6), Inches(3.15), Inches(8.3), Inches(1.3))
tf = paper2_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: Intelligent wearable for real-time seizure detection"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: Combines GSR + Accelerometer with ML algorithms"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Limitation: False positives from physiological changes unrelated to seizures"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['secondary']

# Paper 3: Chatzichristos et al. (SeizeIT2)
add_rounded_box(slide3, 0.3, 4.55, 4.2, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper3_info = slide3.shapes.add_textbox(Inches(0.4), Inches(4.6), Inches(4), Inches(1.3))
tf = paper3_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "3. Chatzichristos, Swinnen, et al."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"SeizeIT2: Multi-modal seizure detection using wearables"'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "OpenNeuro ds005873, Year: 2024"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3, 4.5, 4.55, 8.5, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper3_details = slide3.shapes.add_textbox(Inches(4.6), Inches(4.6), Inches(8.3), Inches(1.3))
tf = paper3_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: ACC, GYR, ECG data from 125 epilepsy patients (883 seizures)"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: Large-scale dataset for ML training with annotated seizures"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Used: Primary dataset for our accelerometer + gyroscope ML model"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['mint']

# Paper 4: Beniczky (ACC detection)
add_rounded_box(slide3, 0.3, 6, 4.2, 1.25, COLORS['light_bg'], "", 10, COLORS['dark'])
paper4_info = slide3.shapes.add_textbox(Inches(0.4), Inches(6.05), Inches(4), Inches(1.2))
tf = paper4_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "4. Beniczky, Polber, Kjaer, et al."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"Detection of tonic-clonic seizures by accelerometry"'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Epilepsia Journal, Year: 2013"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3, 4.5, 6, 8.5, 1.25, COLORS['light_bg'], "", 10, COLORS['dark'])
paper4_details = slide3.shapes.add_textbox(Inches(4.6), Inches(6.05), Inches(8.3), Inches(1.2))
tf = paper4_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: Arm-worn accelerometer for tonic-clonic seizure detection"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: >90% sensitivity for motor seizures, deltoid placement optimal"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']

# ============================================================================
# SLIDE 3b: Literature Survey (Part 2)
# ============================================================================
slide3b = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3b, COLORS['light_bg'])
add_title_bar(slide3b, COLORS['secondary'])

add_text(slide3b, "Literature Survey (Cont.)", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['secondary'], align='center')

# Table header
add_rounded_box(slide3b, 0.3, 1.1, 4.2, 0.5, COLORS['dark'], "PAPER INFORMATION", 11)
add_rounded_box(slide3b, 4.5, 1.1, 8.5, 0.5, COLORS['dark'], "FOCUS AREA, KEY CONTRIBUTIONS, LIMITATIONS", 11)

# Paper 5: Jansen & Lagae (HRV)
add_rounded_box(slide3b, 0.3, 1.65, 4.2, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper5_info = slide3b.shapes.add_textbox(Inches(0.4), Inches(1.7), Inches(4), Inches(1.3))
tf = paper5_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "5. Jansen K., Lagae L."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"Cardiac changes in epilepsy" - Heart Rate Variability in Seizures'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Seizure Journal, Year: 2019"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3b, 4.5, 1.65, 8.5, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper5_details = slide3b.shapes.add_textbox(Inches(4.6), Inches(1.7), Inches(8.3), Inches(1.3))
tf = paper5_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: HRV and cardiac changes during epileptic seizures"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: Ictal tachycardia in 80% of seizures; RMSSD drops significantly"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Used: Basis for our HRV heuristic thresholds (HR >30%, RMSSD <50ms)"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['mint']

# Paper 6: WESAD Dataset
add_rounded_box(slide3b, 0.3, 3.1, 4.2, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper6_info = slide3b.shapes.add_textbox(Inches(0.4), Inches(3.15), Inches(4), Inches(1.3))
tf = paper6_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "6. Schmidt, Reber, Schwabedal, et al."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"WESAD: Wearable Stress and Affect Detection Dataset"'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "ACM ICMI, Year: 2018"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3b, 4.5, 3.1, 8.5, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper6_details = slide3b.shapes.add_textbox(Inches(4.6), Inches(3.15), Inches(8.3), Inches(1.3))
tf = paper6_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: Multi-modal dataset with ACC, BVP, EDA, ECG, EMG, TEMP"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: 15 subjects, stress/affect detection using wearable sensors"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Relevance: GSR/EDA patterns for stress detection applicable to seizures"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['mint']

# Paper 7: Onorati (Multi-modal)
add_rounded_box(slide3b, 0.3, 4.55, 4.2, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper7_info = slide3b.shapes.add_textbox(Inches(0.4), Inches(4.6), Inches(4), Inches(1.3))
tf = paper7_info.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "7. Onorati, Regalia, Caborni, et al."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = '"Multicenter clinical assessment of improved wearable seizure detection"'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Epilepsia Journal, Year: 2017"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['accent']

add_rounded_box(slide3b, 4.5, 4.55, 8.5, 1.4, COLORS['light_bg'], "", 10, COLORS['dark'])
paper7_details = slide3b.shapes.add_textbox(Inches(4.6), Inches(4.6), Inches(8.3), Inches(1.3))
tf = paper7_details.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Focus: Multi-modal wearable combining EDA + accelerometer"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Key: Combining modalities improved specificity from 64% to 92%"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "Used: Basis for our multi-sensor fusion approach"
p.font.size = Pt(9)
p.font.color.rgb = COLORS['mint']

# ============================================================================
# SLIDE 4: Objectives
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, COLORS['light_bg'])
add_title_bar(slide4, COLORS['accent'])

add_text(slide4, "Project Objectives", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['accent'], align='center')

objectives = [
    ("1", "Develop ML Model", "Train a 1D CNN on SeizeIT2 dataset using accelerometer and gyroscope data", COLORS['primary']),
    ("2", "Multi-Sensor Fusion", "Integrate motion (ML), HRV, and GSR detection for robust seizure identification", COLORS['secondary']),
    ("3", "Embedded Deployment", "Deploy on ESP32-S3 microcontroller with TensorFlow Lite Micro", COLORS['accent']),
    ("4", "Real-time Alert", "Implement BLE beacon with SMS notification to caregivers via smartphone", COLORS['mint']),
]

for i, (num, title, desc, color) in enumerate(objectives):
    top = 1.3 + i * 1.4
    
    # Number circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(top), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_text(slide4, title, 1.4, top - 0.05, 4, 0.4, font_size=18, bold=True, color=COLORS['dark'])
    add_text(slide4, desc, 1.4, top + 0.4, 11, 0.5, font_size=14, color=COLORS['dark'])

# ============================================================================
# SLIDE 5: Hardware Components
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, COLORS['light_bg'])
add_title_bar(slide5, COLORS['primary'])

add_text(slide5, "Hardware Components", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['primary'], align='center')

# Main controller
add_rounded_box(slide5, 4, 1.2, 5.3, 1, COLORS['dark'], "ESP32-S3 N16R8", 20)
add_text(slide5, "Dual-Core 240MHz  •  16MB Flash  •  8MB PSRAM", 4, 2.3, 5.3, 0.4, font_size=12, color=COLORS['dark'], align='center')

# Sensors
sensors = [
    ("MPU6050", "Accelerometer\n+ Gyroscope", "±2g, ±250°/s | I2C", COLORS['primary']),
    ("MAX30102", "PPG Sensor\n(Heart Rate)", "IR + Red LEDs | I2C", COLORS['secondary']),
    ("Grove GSR", "Galvanic Skin\nResponse", "Analog | ADC GPIO5", COLORS['accent']),
]

for i, (name, desc, specs, color) in enumerate(sensors):
    left = 0.8 + i * 4.2
    add_rounded_box(slide5, left, 3.2, 3.8, 0.7, color, name, 16)
    add_text(slide5, desc, left, 4, 3.8, 0.7, font_size=14, bold=True, color=COLORS['dark'], align='center')
    add_text(slide5, specs, left, 4.7, 3.8, 0.4, font_size=11, color=COLORS['dark'], align='center')

# Communication
add_rounded_box(slide5, 4, 5.5, 5.3, 0.6, COLORS['mint'], "BLE 5.0 (NimBLE Stack)", 14)
add_text(slide5, "Wireless alert to smartphone → SMS to caregiver", 4, 6.2, 5.3, 0.4, font_size=11, color=COLORS['dark'], align='center')

# ============================================================================
# SLIDE 6: Software Stack
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6, COLORS['light_bg'])
add_title_bar(slide6, COLORS['secondary'])

add_text(slide6, "Software Stack", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['secondary'], align='center')

# Left column
add_text(slide6, "Development Tools", 0.8, 1.2, 5.5, 0.4, font_size=20, bold=True, color=COLORS['primary'])
dev_tools = ["ESP-IDF v5.5.1 (Framework)", "TensorFlow / Keras (Model Training)", 
             "Python 3.12 (Preprocessing)", "C/C++ (Firmware)"]
for i, tool in enumerate(dev_tools):
    add_text(slide6, f"•  {tool}", 1, 1.7 + i*0.45, 5.5, 0.4, font_size=14, color=COLORS['dark'])

# Right column
add_text(slide6, "Key Libraries", 7, 1.2, 5.5, 0.4, font_size=20, bold=True, color=COLORS['secondary'])
libraries = ["TensorFlow Lite Micro (INT8)", "NimBLE Stack (BLE)", 
             "ESP-IDF Driver (I2C, ADC)", "NumPy, Pandas (Data)"]
for i, lib in enumerate(libraries):
    add_text(slide6, f"•  {lib}", 7.2, 1.7 + i*0.45, 5.5, 0.4, font_size=14, color=COLORS['dark'])

# ML Model box
add_rounded_box(slide6, 0.5, 4, 12.333, 1.6, COLORS['dark'])
add_text(slide6, "Machine Learning Model", 0.8, 4.15, 11.7, 0.4, font_size=18, bold=True, color=COLORS['white'], align='center')
add_text(slide6, "1D CNN: Input(206×6) → Conv1D(16) → Conv1D(32) → Conv1D(32) → Dense(1, sigmoid)", 
         0.8, 4.6, 11.7, 0.4, font_size=13, color=RGBColor(200, 200, 200), align='center')
add_text(slide6, "Quantized INT8  •  Model Size: ~15KB  •  Inference: ~10ms on ESP32", 
         0.8, 5.05, 11.7, 0.4, font_size=13, color=COLORS['secondary'], align='center')

# ============================================================================
# SLIDE 7: System Flowchart
# ============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7, COLORS['light_bg'])
add_title_bar(slide7, COLORS['accent'])

add_text(slide7, "System Architecture", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['accent'], align='center')

# Add flowchart image
flowchart_path = '/Users/sonuadithya/Documents/seizure_cursor/template-app/report_graphs/00c_ppt_flowchart.png'
if os.path.exists(flowchart_path):
    slide7.shapes.add_picture(flowchart_path, Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.2))
else:
    add_text(slide7, "[Flowchart: report_graphs/00c_ppt_flowchart.png]", 0.5, 3.5, 12.333, 0.5, font_size=18, color=COLORS['dark'], align='center')

# ============================================================================
# SLIDE 8: Datasets Used
# ============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8, COLORS['light_bg'])
add_title_bar(slide8, COLORS['mint'])

add_text(slide8, "Datasets Used", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['mint'], align='center')

# SeizeIT2 Dataset
add_rounded_box(slide8, 0.3, 1.1, 6.3, 0.5, COLORS['primary'], "SeizeIT2 Dataset (Primary - Motion ML)", 12)
add_rounded_box(slide8, 0.3, 1.65, 6.3, 2.5, COLORS['light_bg'], "", 10, COLORS['dark'])
seizeit_box = slide8.shapes.add_textbox(Inches(0.4), Inches(1.7), Inches(6.1), Inches(2.4))
tf = seizeit_box.text_frame
tf.word_wrap = True
details = [
    "• Patients: 125 | Seizures: 883",
    "• Motor: 80% | Non-Motor: 20%",
    "• Sensors: ACC (25Hz), GYR (25Hz), ECG (256Hz)",
    "• Placement: Upper arm (deltoid muscle)",
    "• Types: Tonic-clonic, Hyperkinetic, Automatisms",
    "• Used for: Accelerometer + Gyroscope ML model"
]
for i, detail in enumerate(details):
    if i == 0:
        tf.paragraphs[0].text = detail
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
    else:
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']

# WESAD Dataset
add_rounded_box(slide8, 6.8, 1.1, 6.2, 0.5, COLORS['secondary'], "WESAD Dataset (Reference - GSR)", 12)
add_rounded_box(slide8, 6.8, 1.65, 6.2, 2.5, COLORS['light_bg'], "", 10, COLORS['dark'])
wesad_box = slide8.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(6), Inches(2.4))
tf = wesad_box.text_frame
tf.word_wrap = True
wesad_details = [
    "• Subjects: 15 | Modalities: 6",
    "• Sensors: ACC, BVP, EDA, ECG, EMG, TEMP",
    "• Focus: Stress and affect detection",
    "• EDA/GSR patterns for stress response",
    "• Sampling: Chest (700Hz), Wrist (varied)",
    "• Used for: GSR heuristic threshold reference"
]
for i, detail in enumerate(wesad_details):
    if i == 0:
        tf.paragraphs[0].text = detail
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
    else:
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']

# Preprocessing summary
add_rounded_box(slide8, 0.3, 4.3, 12.7, 0.5, COLORS['accent'], "Preprocessing Pipeline", 12)
add_rounded_box(slide8, 0.3, 4.85, 12.7, 1.9, COLORS['light_bg'], "", 10, COLORS['dark'])
preprocess_box = slide8.shapes.add_textbox(Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.7))
tf = preprocess_box.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "SeizeIT2 Processing:"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLORS['primary']
p = tf.add_paragraph()
p.text = "Resampled to 20Hz → 10.3s windows (206 samples) → Z-score normalization (μ=-0.048, σ=0.195) → INT8 quantization"
p.font.size = Pt(11)
p.font.color.rgb = COLORS['dark']
p = tf.add_paragraph()
p.text = "HRV Processing:"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS['secondary']
p = tf.add_paragraph()
p.text = "ECG literature → PPG adaptation → Peak detection → R-R intervals → HR, RMSSD, SDNN calculation"
p.font.size = Pt(11)
p.font.color.rgb = COLORS['dark']

# ============================================================================
# SLIDE 9: Methodology
# ============================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide9, COLORS['light_bg'])
add_title_bar(slide9, COLORS['primary'])

add_text(slide9, "Detection Methodology", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['primary'], align='center')

# Three columns
methods = [
    ("Motion Detection", "Machine Learning", ["1D CNN classifier", "Input: 206×6 (ACC+GYR)", "TensorFlow Lite Micro", "INT8 quantized"], "60-70%", COLORS['primary']),
    ("HRV Detection", "Heuristic Rules", ["Peak detection algorithm", "R-R interval calculation", "HR, RMSSD metrics", "30s baseline"], "20-25%", COLORS['secondary']),
    ("GSR Detection", "Heuristic Rules", ["ADC sampling @ 10Hz", "EMA smoothing", "Baseline tracking", ">20% = stress"], "10-15%", COLORS['accent']),
]

for i, (title, subtitle, bullets, weight, color) in enumerate(methods):
    left = 0.5 + i * 4.2
    add_rounded_box(slide9, left, 1.2, 4, 0.55, color, title, 14)
    add_text(slide9, subtitle, left, 1.8, 4, 0.3, font_size=11, color=color, align='center')
    
    for j, bullet in enumerate(bullets):
        add_text(slide9, f"• {bullet}", left + 0.2, 2.2 + j*0.4, 3.6, 0.35, font_size=12, color=COLORS['dark'])
    
    add_text(slide9, f"Weight: {weight}", left, 4, 4, 0.3, font_size=12, bold=True, color=color, align='center')

# Decision logic
add_rounded_box(slide9, 0.5, 4.6, 12.333, 1.2, COLORS['dark'])
add_text(slide9, "Multi-Sensor Fusion Decision Rules", 0.8, 4.75, 11.7, 0.3, font_size=16, bold=True, color=COLORS['white'], align='center')
add_text(slide9, "Motion > 50% → ALERT   |   2+ Sensors Agree → ALERT   |   Motion > 30% + Support → ALERT", 
         0.8, 5.2, 11.7, 0.4, font_size=14, color=COLORS['warm'], align='center')

# ============================================================================
# SLIDE 10: Results - Performance
# ============================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide10, COLORS['light_bg'])
add_title_bar(slide10, COLORS['secondary'])

add_text(slide10, "Model Performance", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['secondary'], align='center')

# Metrics
metrics = [("92.4%", "Sensitivity", COLORS['primary']), ("95.7%", "Specificity", COLORS['secondary']),
           ("87.6%", "Precision", COLORS['accent']), ("94.8%", "Accuracy", COLORS['mint'])]
for i, (value, metric, color) in enumerate(metrics):
    left = 0.7 + i * 3.2
    add_text(slide10, value, left, 1.2, 2.8, 0.8, font_size=40, bold=True, color=color, align='center')
    add_text(slide10, metric, left, 2, 2.8, 0.4, font_size=18, bold=True, color=COLORS['dark'], align='center')

# Confusion matrix image
confusion_path = '/Users/sonuadithya/Documents/seizure_cursor/template-app/report_graphs/07_confusion_matrix.png'
if os.path.exists(confusion_path):
    slide10.shapes.add_picture(confusion_path, Inches(0.5), Inches(2.7), Inches(12.333), Inches(4.5))

# ============================================================================
# SLIDE 11: Training Curves
# ============================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide11, COLORS['light_bg'])
add_title_bar(slide11, COLORS['accent'])

add_text(slide11, "Training Performance", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['accent'], align='center')

curves_path = '/Users/sonuadithya/Documents/seizure_cursor/template-app/report_graphs/06_training_curves.png'
if os.path.exists(curves_path):
    slide11.shapes.add_picture(curves_path, Inches(0.5), Inches(1.1), Inches(12.333), Inches(5.2))

add_text(slide11, "✓ Converges in ~30 epochs   •   ✓ No overfitting   •   ✓ Validation accuracy ~92%", 
         0.5, 6.4, 12.333, 0.4, font_size=14, color=COLORS['dark'], align='center')

# ============================================================================
# SLIDE 12: Signal Comparison
# ============================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide12, COLORS['light_bg'])
add_title_bar(slide12, COLORS['mint'])

add_text(slide12, "Seizure vs Normal Signals", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['mint'], align='center')

signal_path = '/Users/sonuadithya/Documents/seizure_cursor/template-app/report_graphs/09_signal_comparison.png'
if os.path.exists(signal_path):
    slide12.shapes.add_picture(signal_path, Inches(0.2), Inches(1), Inches(12.9), Inches(6.3))

# ============================================================================
# SLIDE 13: Observations
# ============================================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide13, COLORS['light_bg'])
add_title_bar(slide13, COLORS['primary'])

add_text(slide13, "System Observations", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['primary'], align='center')

observations = [
    ("Real-time Performance", "System maintains 50ms loop rate with all three sensors active"),
    ("Detection Latency", "Seizure alert triggered within 2 seconds of onset"),
    ("Sensor Reliability", "MPU6050 and MAX30102 provide consistent readings via I2C"),
    ("BLE Alert", "Beacon detected by smartphone within 5 meters range"),
    ("Power Consumption", "~150mA during monitoring (suitable for battery operation)"),
    ("False Positive Rate", "< 5% during normal activities (walking, sitting, typing)"),
]

for i, (title, desc) in enumerate(observations):
    top = 1.15 + i * 0.95
    add_text(slide13, f"✓  {title}:", 0.8, top, 4.5, 0.35, font_size=15, bold=True, color=COLORS['primary'])
    add_text(slide13, desc, 1, top + 0.35, 11.5, 0.4, font_size=13, color=COLORS['dark'])

# ============================================================================
# SLIDE 14: Future Work
# ============================================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide14, COLORS['light_bg'])
add_title_bar(slide14, COLORS['secondary'])

add_text(slide14, "Future Enhancements", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['secondary'], align='center')

future = [
    ("Cloud Integration", "Send seizure data to cloud for remote monitoring", COLORS['primary']),
    ("GPS Location", "Include location in SMS for emergency response", COLORS['secondary']),
    ("Fall Detection", "Detect falls during or after seizures", COLORS['accent']),
    ("Battery Optimization", "Sleep modes for 24+ hour battery life", COLORS['mint']),
    ("Mobile App", "Seizure logging, trends, medication reminders", COLORS['sky']),
    ("EEG Integration", "Dry EEG for non-motor seizure detection", COLORS['warm']),
]

for i, (title, desc, color) in enumerate(future):
    row, col = i // 2, i % 2
    left = 0.5 + col * 6.5
    top = 1.2 + row * 1.8
    add_rounded_box(slide14, left, top, 6, 0.55, color, title, 14)
    add_text(slide14, desc, left + 0.1, top + 0.65, 5.8, 0.6, font_size=12, color=COLORS['dark'])

# ============================================================================
# SLIDE 15: Conclusion
# ============================================================================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide15, COLORS['light_bg'])
add_title_bar(slide15, COLORS['primary'])

add_text(slide15, "Conclusion", 0.5, 0.4, 12.333, 0.6, font_size=36, bold=True, color=COLORS['primary'], align='center')

conclusions = [
    "Successfully developed a multi-sensor seizure detection wearable",
    "Achieved 92.4% sensitivity and 95.7% specificity",
    "Real-time detection with < 2 second latency",
    "Cost-effective solution using off-the-shelf components",
    "Potential to improve quality of life for epilepsy patients"
]

for i, conclusion in enumerate(conclusions):
    add_text(slide15, f"✓  {conclusion}", 1.5, 1.4 + i*0.6, 10.5, 0.5, font_size=17, color=COLORS['dark'])

# Thank you
add_text(slide15, "Thank You!", 0.5, 5, 12.333, 0.7, font_size=40, bold=True, color=COLORS['secondary'], align='center')
add_text(slide15, "Questions?", 0.5, 5.8, 12.333, 0.5, font_size=22, color=COLORS['dark'], align='center')

# ============================================================================
# Save
# ============================================================================
output_path = '/Users/sonuadithya/Documents/seizure_cursor/template-app/Seizure_Detection_Presentation.pptx'
prs.save(output_path)

print("✅ PowerPoint Presentation Generated!")
print(f"   Location: {output_path}")
print("   Slides: 15")
print("   Design: Clean, minimalistic, no overlapping elements")
